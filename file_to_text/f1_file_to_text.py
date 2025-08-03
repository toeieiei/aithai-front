from PIL import Image
import fitz
from pptx import Presentation
import pdfplumber
from fixthaipdf import clean
from pptx import Presentation
import re
import os
import csv
import glob
import pandas as pd
import requests
import easyocr
# from dotenv import load_dotenv
from docx import Document
from docx2pdf import convert
import aspose.words as aw
import aspose.slides as slides


# load_dotenv()
# API_KEY = os.getenv("OCR_API_KEY")

# from ocr_model import get_easyocr_reader



reader = easyocr.Reader(['en', 'th'], gpu=True)

def extract_text_from_pdf(image_path):
    # extract text from pdf with pdfplumber
    with pdfplumber.open(image_path) as pdf:
        text = ''
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + '\n'
        return text

def extract_text_from_pptx(pptx_path):
    # extract text from pptx with python-pptx
    prs = Presentation(pptx_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text
# def extract_text_from_image(image_path):

#     url = "https://api.aiforthai.in.th/ocr"

#     files = {'uploadfile':open(image_path, 'rb')}

#     headers = {
#         'Apikey': "TVYtq6nCf2hW6BlmSXEBrISUCeyrejnw",
#         }
#     response = requests.post(url, files=files, headers=headers)

    # return response.json()
def extract_text_from_image_easyocr(image_path):
    result = reader.readtext(image_path, detail=0, paragraph=True)
    return result

def extract_text_docx(docx_path):
    doc = Document(docx_path)
    text = ''
    for para in doc.paragraphs:
        text += para.text + '\n'
    return text

def extract_text_from_pdf_images(pdf_path):
    doc = fitz.open(pdf_path)
    text = ''
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=300)  # You can adjust DPI
        img_path = f"temp_pdf_page_{i}.png"
        pix.save(img_path)

        text += extract_text_from_image_easyocr(img_path) + "\n"
        os.remove(img_path)

    return text


def extract_text_from_docx_images(docx_path, image_output_dir="docx_images"):
    # Load the DOCX
    doc = aw.Document(docx_path)
    os.makedirs(image_output_dir, exist_ok=True)

    # Save options
    save_options = aw.saving.ImageSaveOptions(aw.SaveFormat.JPEG)

    full_text = ""
    for page in range(doc.page_count):
        save_options.page_set = aw.saving.PageSet(page)
        image_path = os.path.join(image_output_dir, f"{os.path.basename(docx_path)}_page_{page + 1}.jpg")
        doc.save(image_path, save_options)

        # OCR step
        ocr_result = extract_text_from_image_easyocr(image_path)

        # Make sure it's a string
        if isinstance(ocr_result, list):
            ocr_text = "\n".join(ocr_result)
        else:
            ocr_text = str(ocr_result)

        full_text += ocr_text + "\n"

    return full_text


def extract_text_from_pptx_images(pptx_path, output_dir="pptx_images", format="jpeg"):
    os.makedirs(output_dir, exist_ok=True)
    full_text = ""

    with slides.Presentation(pptx_path) as presentation:
        for i, slide in enumerate(presentation.slides):
            image_path = os.path.join(output_dir, f"slide_{i+1}.{format}")
            
            # Save slide as image
            slide.get_thumbnail(2, 2).save(image_path, getattr(slides.export.SaveFormat, format.upper()))

            # OCR for each image
            ocr_result = extract_text_from_image_easyocr(image_path)

            # Ensure text is a string
            if isinstance(ocr_result, list):
                ocr_text = "\n".join(ocr_result)
            else:
                ocr_text = str(ocr_result)

            full_text += ocr_text + "\n"

    return full_text